"""Board pack export (VA-P7-08): HTML with cover/ToC/sections, PDF, PPTX; branding applied."""

from __future__ import annotations

import html as _html
import re
from io import BytesIO
from typing import Any

try:
    from xhtml2pdf import pisa
except ImportError:
    pisa = None  # type: ignore[assignment]

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None  # type: ignore[assignment, misc]


def _format_table(rows: list[dict], period_labels: list[str] | None = None) -> str:
    """Render list of period dicts as HTML table."""
    if not rows:
        return "<p>No data</p>"
    keys = [k for k in rows[0].keys() if k != "period_index"]
    n_periods = len(rows)
    labels = period_labels or [f"P{i}" for i in range(n_periods)]
    html = ['<table class="pack-table"><thead><tr><th>Line Item</th>']
    for lbl in labels:
        html.append(f"<th>{_html.escape(str(lbl))}</th>")
    html.append("</tr></thead><tbody>")
    for key in keys:
        html.append(f"<tr><td>{_html.escape(key.replace('_', ' ').title())}</td>")
        for i, row in enumerate(rows):
            val = row.get(key)
            if isinstance(val, (int, float)):
                html.append(f"<td>{val:,.2f}</td>")
            else:
                html.append(f"<td>{_html.escape(str(val) if val is not None else '')}</td>")
        html.append("</tr>")
    html.append("</tbody></table>")
    return "".join(html)


def build_board_pack_html(
    label: str,
    section_order: list[str],
    narrative: dict[str, Any],
    statements: dict[str, Any],
    kpis: list[dict[str, Any]],
    budget_summary: str | None,
    branding: dict[str, Any],
    run_id: str = "",
    display_currency: str | None = None,
    benchmark_metrics: list[dict[str, Any]] | None = None,
) -> str:
    """Build full HTML document: cover, ToC, sections in order, branding. Optional display_currency (VA-P8-01), benchmark_metrics (VA-P8-09)."""
    _raw_color = (branding.get("primary_color") or "#2563eb").strip()
    primary_color = _raw_color if re.fullmatch(r"#[0-9a-fA-F]{3,8}", _raw_color) else "#2563eb"
    terms_footer = _html.escape(str(branding.get("terms_footer") or ""))
    _raw_logo = branding.get("logo_url") or ""
    logo_url = _raw_logo if _raw_logo.startswith("https://") else ""

    is_list = statements.get("income_statement") or []
    bs_list = statements.get("balance_sheet") or []
    cf_list = statements.get("cash_flow") or []
    periods = statements.get("periods") or [f"P{i}" for i in range(max(len(is_list), len(bs_list), len(cf_list), 1))]

    section_contents: dict[str, str] = {}
    section_contents["executive_summary"] = f"<div class='section-body'><p>{_html.escape(narrative.get('executive_summary') or 'No executive summary generated.')}</p></div>"
    section_contents["income_statement"] = f"<h3>Income Statement</h3>{_format_table(is_list, periods[: len(is_list)])}"
    section_contents["balance_sheet"] = f"<h3>Balance Sheet</h3>{_format_table(bs_list, periods[: len(bs_list)])}"
    section_contents["cash_flow"] = f"<h3>Cash Flow</h3>{_format_table(cf_list, periods[: len(cf_list)])}"
    section_contents["budget_variance"] = f"<h3>Budget Variance</h3><pre class='variance-pre'>{_html.escape(budget_summary or 'No budget variance data.')}</pre>"
    section_contents["kpi_dashboard"] = f"<h3>KPI Dashboard</h3>{_format_table(kpis, periods[: len(kpis)]) if kpis else '<p>No KPI data.</p>'}"
    section_contents["scenario_comparison"] = "<h3>Scenario Comparison</h3><p>Scenario comparison data can be included when scenario runs are linked.</p>"
    section_contents["strategic_commentary"] = f"<div class='section-body'><p>{_html.escape(narrative.get('strategic_commentary') or 'No strategic commentary generated.')}</p></div>"
    if benchmark_metrics:
        bench_rows = "".join(
            f"<tr><td>{_html.escape(str(m.get('metric_name', '')))}</td><td>{m.get('median', '')}</td><td>{m.get('sample_count', '')} peers</td></tr>"
            for m in benchmark_metrics
        )
        section_contents["benchmark"] = f"<h3>Industry Benchmark (Peer Median)</h3><table class='pack-table'><thead><tr><th>Metric</th><th>Peer Median</th><th>Sample</th></tr></thead><tbody>{bench_rows}</tbody></table>"

    toc_items = []
    body_sections = []
    for i, sec_key in enumerate(section_order, start=1):
        if sec_key not in section_contents:
            continue
        safe_title = _html.escape(sec_key.replace("_", " ").title())
        toc_items.append(f"<li><a href='#sec-{i}'>{safe_title}</a></li>")
        body_sections.append(f"<section id='sec-{i}' class='pack-section'><h2>{safe_title}</h2>{section_contents[sec_key]}</section>")

    logo_img = f"<img src='{_html.escape(logo_url)}' alt='Logo' class='pack-logo'/>" if logo_url else ""
    currency_note = f"<p class='pack-meta'>Amounts in {_html.escape(display_currency)}</p>" if display_currency else ""
    cover = f"""
    <div class="pack-cover" style="border-top:4px solid {primary_color};">
        {logo_img}
        <h1 class="pack-title">{_html.escape(label)}</h1>
        <p class="pack-meta">Run: {_html.escape(run_id or 'N/A')}</p>
        {currency_note}
    </div>
    """
    toc = f"<nav class='pack-toc'><h2>Contents</h2><ol>{''.join(toc_items)}</ol></nav>"
    body_html = cover + toc + "".join(body_sections)
    if terms_footer:
        body_html += f"<footer class='pack-footer' style='border-top:2px solid {primary_color};'><p>{terms_footer}</p></footer>"

    css = f"""
    body {{ font-family: system-ui, sans-serif; margin: 2rem; max-width: 900px; color: #1a1a1a; }}
    .pack-cover {{ padding: 2rem 0; margin-bottom: 2rem; }}
    .pack-logo {{ max-height: 60px; margin-bottom: 1rem; }}
    .pack-title {{ font-size: 1.75rem; color: {primary_color}; margin: 0.5rem 0; }}
    .pack-meta {{ color: #666; font-size: 0.9rem; }}
    .pack-toc {{ margin: 2rem 0; padding: 1rem; background: #f8fafc; border-radius: 8px; }}
    .pack-toc ol {{ margin: 0.5rem 0 0 1.5rem; }}
    .pack-toc a {{ color: {primary_color}; text-decoration: none; }}
    .pack-section {{ margin: 2rem 0; page-break-inside: avoid; }}
    .pack-section h2 {{ font-size: 1.25rem; border-bottom: 1px solid #e2e8f0; padding-bottom: 0.5rem; color: {primary_color}; }}
    .pack-table {{ border-collapse: collapse; width: 100%; margin: 1rem 0; }}
    .pack-table th, .pack-table td {{ border: 1px solid #e2e8f0; padding: 0.5rem; text-align: right; }}
    .pack-table th:first-child, .pack-table td:first-child {{ text-align: left; }}
    .pack-table th {{ background: #f1f5f9; }}
    .variance-pre {{ white-space: pre-wrap; font-size: 0.9rem; background: #f8fafc; padding: 1rem; border-radius: 4px; }}
    .pack-footer {{ margin-top: 3rem; padding-top: 1rem; font-size: 0.8rem; color: #64748b; }}
    """

    return f"""<!DOCTYPE html><html><head><meta charset='utf-8'/><title>{_html.escape(label)}</title><style>{css}</style></head><body>{body_html}</body></html>"""


def html_to_pdf(html: str) -> bytes:
    """Render HTML to PDF bytes. Requires xhtml2pdf."""
    if pisa is None:
        raise RuntimeError("PDF export requires xhtml2pdf; pip install xhtml2pdf")
    buf = BytesIO()
    status = pisa.CreatePDF(html.encode("utf-8"), dest=buf, encoding="utf-8")
    if status.err:
        raise RuntimeError("PDF generation failed")
    return buf.getvalue()


def build_board_pack_pptx(
    label: str,
    section_order: list[str],
    narrative: dict[str, Any],
    statements: dict[str, Any],
    kpis: list[dict[str, Any]],
    budget_summary: str | None,
    run_id: str = "",
    display_currency: str | None = None,
    benchmark_metrics: list[dict[str, Any]] | None = None,
) -> bytes:
    """Build PPTX: title slide + one slide per section (title + content summary). VA-P7-08; VA-P8-01 display_currency."""
    if Presentation is None:
        raise RuntimeError("PPTX export requires python-pptx; pip install python-pptx")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(32)
    p.font.bold = True
    run_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.5))
    run_box.text_frame.text = f"Run: {run_id or 'N/A'}" + (f"  |  Amounts in {display_currency}" if display_currency else "")

    section_contents: dict[str, tuple[str, str]] = {}
    section_contents["executive_summary"] = ("Executive Summary", narrative.get("executive_summary") or "—")
    section_contents["income_statement"] = ("Income Statement", "See appendix for full table.")
    section_contents["balance_sheet"] = ("Balance Sheet", "See appendix for full table.")
    section_contents["cash_flow"] = ("Cash Flow", "See appendix for full table.")
    section_contents["budget_variance"] = ("Budget Variance", budget_summary or "No variance data.")
    section_contents["kpi_dashboard"] = ("KPI Dashboard", "See appendix for full KPIs." if kpis else "No KPI data.")
    section_contents["scenario_comparison"] = ("Scenario Comparison", "Scenario data when linked.")
    section_contents["strategic_commentary"] = ("Strategic Commentary", narrative.get("strategic_commentary") or "—")
    if benchmark_metrics:
        bench_text = "; ".join(f"{m.get('metric_name', '')}: median {m.get('median', '')}" for m in benchmark_metrics[:10])
        section_contents["benchmark"] = ("Industry Benchmark", bench_text or "No benchmark data.")

    for sec_key in section_order:
        if sec_key not in section_contents:
            continue
        title, content = section_contents[sec_key]
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75)).text_frame.text = title
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        tf = body_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content[:3000] + ("..." if len(content) > 3000 else "")
        p.font.size = Pt(14)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
